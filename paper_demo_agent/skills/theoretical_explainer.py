"""Skill for theory/math papers — adapts to app / presentation / website / slides / latex."""

from paper_demo_agent.paper.models import Paper, PaperAnalysis
from paper_demo_agent.skills.base import BaseSkill, FORM_SPECS


# ─────────────────────────────────────────────────────────────────────────────
# Pre-baked tool knowledge — never search for these basics
# ─────────────────────────────────────────────────────────────────────────────

_PYTHON_PPTX_PATTERNS = """
━━ PYTHON-PPTX 1.0.0 REFERENCE (use verbatim, do NOT search) ━━

SETUP:
```python
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import os

prs = Presentation()
prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)  # 16:9
BLANK = prs.slide_layouts[6]  # always use blank layout
```

HELPERS (copy these into build.py):
```python
def set_dark_bg(slide, color=RGBColor(0x09, 0x09, 0x0b)):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color

def add_header_bar(slide, title_text, accent=RGBColor(0x63, 0x66, 0xf1)):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    p = bar.text_frame.paragraphs[0]
    p.text = title_text; p.font.size = Pt(28); p.font.bold = True
    p.font.color.rgb = RGBColor(0xff,0xff,0xff); p.alignment = PP_ALIGN.LEFT
    bar.text_frame.margin_left = Inches(0.3); bar.text_frame.margin_top = Inches(0.15)

def add_text_box(slide, text, left, top, width, height, size=Pt(20), bold=False,
                 color=RGBColor(0xfa,0xfa,0xfa), align=PP_ALIGN.LEFT):
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = True; p = tf.paragraphs[0]
    p.text = text; p.font.size = size; p.font.bold = bold
    p.font.color.rgb = color; p.alignment = align

def add_bullet_list(slide, items, left, top, width, height, size=Pt(20)):
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"• {item}"; p.font.size = size
        p.font.color.rgb = RGBColor(0xfa,0xfa,0xfa); p.space_before = Pt(6)
```

CHART (CRITICAL — do NOT use .style property, set colors via format.fill):
```python
chart_data = ChartData()
chart_data.categories = ['Ours', 'Baseline A', 'Baseline B']
chart_data.add_series('Score', (45.3, 38.1, 41.7))
chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1), Inches(1.2), Inches(11), Inches(5.5), chart_data).chart
chart.has_title = True
chart.chart_title.text_frame.text = 'Main Benchmark Results'
chart.has_legend = True
chart.chart_area.format.fill.solid(); chart.chart_area.format.fill.fore_color.rgb = RGBColor(0x09,0x09,0x0b)
chart.plot_area.format.fill.solid(); chart.plot_area.format.fill.fore_color.rgb = RGBColor(0x11,0x11,0x13)
chart.value_axis.has_major_gridlines = True
for i, series in enumerate(chart.plots[0].series):
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = [RGBColor(0x63,0x66,0xf1), RGBColor(0x22,0xc5,0x5e)][i % 2]
chart.chart_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xfa,0xfa,0xfa)
chart.category_axis.tick_labels.font.color.rgb = RGBColor(0xfa,0xfa,0xfa)
chart.value_axis.tick_labels.font.color.rgb = RGBColor(0xfa,0xfa,0xfa)
chart.legend.font.color.rgb = RGBColor(0xfa,0xfa,0xfa)
chart.plots[0].has_data_labels = True
chart.plots[0].data_labels.show_value = True
```

TABLE (styled results comparison):
```python
def add_results_table(slide, headers, rows, highlight_row=0):
    tbl = slide.shapes.add_table(len(rows)+1, len(headers),
        Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.2)).table
    ACCENT = RGBColor(0x63,0x66,0xf1); WHITE = RGBColor(0xff,0xff,0xff)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = ACCENT
        c.text_frame.paragraphs[0].font.bold = True; c.text_frame.paragraphs[0].font.color.rgb = WHITE
    for i, row in enumerate(rows):
        is_ours = (i == highlight_row)
        bg = RGBColor(0x1e,0x1b,0x4b) if is_ours else RGBColor(0x18,0x18,0x1b)
        for j, val in enumerate(row):
            c = tbl.cell(i+1, j); c.text = str(val); c.fill.solid(); c.fill.fore_color.rgb = bg
            c.text_frame.paragraphs[0].font.bold = is_ours
            c.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x22,0xc5,0x5e) if is_ours else RGBColor(0xfa,0xfa,0xfa)
    if len(headers) > 0: tbl.columns[0].width = Inches(2.8)
```

MATPLOTLIB DIAGRAM → IMAGE:
```python
import matplotlib; matplotlib.use('Agg')
import matplotlib.pyplot as plt; from io import BytesIO
fig, ax = plt.subplots(figsize=(10, 4)); fig.patch.set_facecolor('#09090b'); ax.set_facecolor('#111113')
# ... draw boxes, arrows, text ...
ax.axis('off'); buf = BytesIO()
plt.savefig(buf, format='png', bbox_inches='tight', dpi=300); buf.seek(0); plt.close()
slide.shapes.add_picture(buf, Inches(1.5), Inches(1.2), Inches(10.0), Inches(5.5))
```

SPEAKER NOTES: `slide.notes_slide.notes_text_frame.text = "Notes here"`
IMAGE FROM FILE: `slide.shapes.add_picture('figures/fig1.png', Inches(1.5), Inches(1.2), Inches(10), Inches(5.5))`
SAVE: `prs.save('presentation.pptx')`

━━ CRITICAL RULES ━━
• NEVER use pdfplumber, fitz, or any PDF parser to extract table data at runtime.
  All table data must be manually read from the paper and hard-coded as Python lists.
• NEVER dump raw PDF text into single-column tables. Tables must be properly structured
  with typed headers and data columns (e.g., ['Method', 'BLEU', 'ROUGE']).
• For dark-theme slides, charts MUST style title/legend/axes/ticks to light text and include data labels
  for key result values. Default python-pptx chart colors are often unreadable.
• Include only 2-3 key results tables. Skip hyperparameter configs, training details, appendix tables.
• 16-20 slides maximum. Every slide must advance the narrative.
"""

_LATEX_BEAMER_PATTERNS = """
━━ LATEX BEAMER REFERENCE (use verbatim, do NOT search) ━━

PREAMBLE:
```latex
\\documentclass[aspectratio=169,11pt]{beamer}
\\usetheme{metropolis}
\\definecolor{accent}{HTML}{6366f1}
\\definecolor{bgdark}{HTML}{09090b}
\\definecolor{textlight}{HTML}{fafafa}
\\setbeamercolor{frametitle}{fg=white, bg=accent}
\\setbeamercolor{normal text}{fg=textlight, bg=bgdark}
\\setbeamercolor{alerted text}{fg=accent}
\\setbeamercolor{block title}{fg=white, bg=accent}
\\setbeamercolor{block body}{fg=textlight, bg=bgdark!80}
\\usepackage{amsmath, amssymb, amsthm, tikz, booktabs, graphicx, hyperref, listings}
\\usetikzlibrary{arrows.meta, positioning, fit, calc, shapes.geometric}
\\lstset{basicstyle=\\ttfamily\\small\\color{textlight}, backgroundcolor=\\color{bgdark!80},
  breaklines=true, frame=single, rulecolor=\\color{accent}}
\\title{Paper Title} \\subtitle{Venue} \\author{Authors} \\date{\\today}
```

KEY PATTERNS:
- Title: `\\begin{document} \\maketitle`
- Bullets: `\\item<1-> First point` (incremental reveals)
- Math: `\\begin{equation} ... \\end{equation}` or `\\begin{align} ... \\end{align}`
- Theorem: `\\begin{theorem}[Name] ... \\end{theorem} \\pause \\begin{block}{Intuition} ... \\end{block}`
- Table: `\\begin{tabular}{lcc} \\toprule ... \\midrule ... \\bottomrule \\end{tabular}`
  Use `\\textbf{}` to highlight best results
- TikZ: `\\begin{tikzpicture}[node distance=1.8cm, box/.style={draw, rounded corners, fill=accent!30, text=white}]`
  `\\node[box] (a) {Input}; \\node[box, right=of a] (b) {Encoder}; \\draw[->, thick] (a) -- (b);`
- Columns: `\\begin{columns}[T] \\begin{column}{0.48\\textwidth} ... \\end{column} ... \\end{columns}`
- Figure: `\\includegraphics[width=0.88\\textwidth]{figures/fig1.png}`
- End: `\\begin{frame}{Thank You} \\centering \\Large Questions? \\end{frame} \\end{document}`
"""


class TheoreticalExplainerSkill(BaseSkill):
    name = "TheoreticalExplainerSkill"
    description = "Theory/math paper → slides (pptx), latex (beamer), presentation (reveal.js), or website"

    def get_system_prompt(self, paper: Paper, analysis: PaperAnalysis, demo_form: str, demo_type: str) -> str:
        # Inject form-specific pre-baked patterns
        form_knowledge = ""
        if demo_form == "slides":
            form_knowledge = _PYTHON_PPTX_PATTERNS
        elif demo_form == "latex":
            form_knowledge = _LATEX_BEAMER_PATTERNS
        elif demo_form == "presentation":
            form_knowledge = self._reveal_patterns()

        return f"""You are a world-class science communicator combining the depth of a MIT
professor, the clarity of 3Blue1Brown, and the visual craft of a NeurIPS keynote speaker.
You make abstract mathematics viscerally intuitive without dumbing it down.

PAPER CONTEXT:
{self._paper_summary(paper, analysis)}

{self._form_block(demo_form)}

{self._demo_type_guidance(demo_type)}

━━ SKILL CONTEXT — Theory / Mathematics Paper ━━

STEP 0 — DECODE THE THEORETICAL CONTRIBUTION
From the paper, identify:
  • What is the CORE CLAIM? (one sentence theorem / impossibility result / bound)
  • What does the reader need to understand BEFORE they can appreciate the proof?
  • What is the KEY INSIGHT that makes the proof work?
  • What are the PRACTICAL IMPLICATIONS of this result?
  • What is a concrete NUMERICAL EXAMPLE that illustrates the theorem?

EXPLANATION FRAMEWORK — apply to EVERY major concept:
  1. 🔥 Motivation  → "What breaks without this? Here's a concrete failure case."
  2. 💡 Intuition   → "The big idea in one sentence, using an analogy."
  3. 📐 Formalism   → "The precise definition/theorem, typeset with math."
  4. 🔢 Example     → "Here's a concrete case with real numbers."
  5. 🌍 Implication → "This matters because in practice it means X."

SLIDE STRUCTURE — build this exact sequence:
  1.  Title — paper title, authors, venue, year
  2.  Problem Statement — what question is this paper answering? (with failure examples)
  3-4. Background — 2 slides of prerequisite concepts
  5.  Key Insight — the "aha moment" (large font, minimal text, strong visual)
  6-8. Method — 3 slides: informal overview → formal statement → proof sketch / derivation
  9.  Concrete Example — step-by-step numerical walkthrough
  10. Main Result — the theorem/bound stated formally
  11. Implications — practical consequences and why this matters
  12. Comparison — how does this relate to / improve prior results?
  13. Limitations & Open Questions
  14. Conclusion — 3 key takeaways (one at a time)
  15. Q&A

{form_knowledge}

WEBSITE FORM PATTERNS (when form=website):
  • Scrollytelling layout: sections enter with IntersectionObserver animations
  • Math with KaTeX auto-render: `renderMathInElement(document.body, {{delimiters:[...]}})`
  • Each concept in a "card" with: emoji icon, title, informal → formal toggle button
  • Interactive parameter slider for key equations

GRADIO APP PATTERNS (when form=app):
  • gr.Markdown with full LaTeX support (Gradio renders $$ out of the box)
  • gr.Slider for theorem parameters → gr.Markdown updates to show how the bound changes
  • gr.Code to show proof steps in syntax-highlighted code
  • Tabs: Intuition | Formal Proof | Interactive Example | Key Results

{self._multistep_instructions(demo_form)}

{self._tool_usage_instructions()}
"""

    def _reveal_patterns(self) -> str:
        return """━━ REVEAL.JS PATTERNS (use verbatim — CDN URLs are pre-baked in form spec) ━━

  # Dark branded theme — override CSS variables in <style> inside <head>:
  ```css
  :root {
    --r-background-color: #09090b;
    --r-main-font: 'Inter', sans-serif;
    --r-main-font-size: 34px;
    --r-main-color: #fafafa;
    --r-heading-font: 'Inter', sans-serif;
    --r-heading-color: #f0f0ff;
    --r-heading-font-weight: 700;
    --r-link-color: #818cf8;
    --r-selection-background-color: rgba(99,102,241,0.4);
  }
  .reveal { background: #09090b; }
  .reveal .slides section { padding: 0 60px; }
  .math-block { background: rgba(99,102,241,0.1); border-left: 3px solid #6366f1;
    padding: 16px 24px; border-radius: 0 8px 8px 0; margin: 16px 0; }
  ```

  # data-auto-animate (connect consecutive slides):
  ```html
  <section data-auto-animate>
    <h3>Step 1: Initialize</h3>
    <p data-id="eq1">$$ x_0 = 0 $$</p>
  </section>
  <section data-auto-animate>
    <h3>Step 2: Iterate</h3>
    <p data-id="eq1">$$ x_{t+1} = x_t - \\eta \\nabla f(x_t) $$</p>
    <p class="fragment">where $\\eta$ is the learning rate</p>
  </section>
  ```

  # Fragments — never show all bullets at once:
  ```html
  <section>
    <h3>Three Key Properties</h3>
    <ul>
      <li class="fragment">Property 1 <span class="fragment highlight-red">(critical!)</span></li>
      <li class="fragment">Property 2</li>
      <li class="fragment">Property 3</li>
    </ul>
  </section>
  ```

  # Inline SVG diagrams — draw at least 3:
  ```html
  <svg width="600" height="200" viewBox="0 0 600 200">
    <defs>
      <marker id="arrow" viewBox="0 0 10 10" refX="10" refY="5"
              markerWidth="6" markerHeight="6" orient="auto">
        <path d="M 0 0 L 10 5 L 0 10 z" fill="#818cf8"/>
      </marker>
    </defs>
    <circle cx="100" cy="100" r="80" stroke="#6366f1" fill="none" stroke-width="2"/>
    <line x1="100" y1="100" x2="160" y2="60" stroke="#22c55e"
          stroke-width="2" marker-end="url(#arrow)"/>
    <text x="165" y="58" fill="#fafafa" font-size="14">x</text>
  </svg>
  ```

  # KaTeX math — ALWAYS use $$ (block) or $ (inline):
  ```html
  <p>The loss is $\\mathcal{L}(x) = -\\log p_\\theta(x)$</p>
  <div class="math-block">
    $$ \\mathbb{E}[f(X)] \\leq \\frac{1}{2} \\sigma^2 L^2 $$
  </div>
  ```

  # Reveal.initialize() — must include ALL plugins:
  ```javascript
  Reveal.initialize({
    hash: true, slideNumber: 'c/t',
    transition: 'slide', transitionSpeed: 'fast',
    backgroundTransition: 'fade',
    autoAnimateEasing: 'ease-out', autoAnimateDuration: 0.6,
    plugins: [ RevealHighlight, RevealMath.KaTeX, RevealNotes ]
  });
  ```
"""

    def get_initial_message(self, paper: Paper, analysis: PaperAnalysis, demo_form: str, demo_type: str) -> str:
        form_hint = {
            "slides":       "python-pptx (.pptx file via build.py)",
            "latex":        "LaTeX/Beamer (.tex file)",
            "presentation": "reveal.js HTML (demo.html)",
            "website":      "static HTML/CSS/JS (index.html)",
            "app":          "Gradio Python app (app.py)",
        }.get(demo_form, demo_form)

        figures_note = ""
        if demo_form == "slides":
            figures_note = """
STEP 0 — BEFORE WRITING ANY SLIDES (do this first):
  a) Extract key figures and tables from the paper PDF — TWO-STEP process:
     Step 1: Extract FULL pages (no crop) to see what each page contains:
       extract_pdf_page(page=3)  → method diagram page
       extract_pdf_page(page=5)  → results page
       extract_pdf_page(page=7)  → table/chart page
     Step 2: Only crop if a figure clearly occupies a small region of the page.
       If unsure, use the FULL page PNG — `add_picture('figures/page_3.png', ...)` scales fine.
       DO NOT invent crop coordinates — guessing y0/y1 will often capture the wrong content.
  b) Parse ALL numeric results from the paper text — write them down before coding
  c) Branding:
     - Pick accent color based on paper institution (e.g., Google blue #4285F4, Meta blue #0082FB, or default #6366f1)
     - Do NOT download institution logos as SVG — python-pptx cannot render SVGs
     - Use text labels or colored shapes for branding instead
  d) install_package: python-pptx (if needed)

"""
        elif demo_form == "latex":
            figures_note = """
STEP 0 — MANDATORY: DO THIS BEFORE WRITING ANY FRAMES
  You MUST extract figures from the paper PDF before writing any LaTeX frames.
  paper.pdf is already in the output directory.

  TWO-STEP EXTRACTION (critical — do NOT guess crop coordinates):
  Step 0a — Extract FULL pages first (no crop arg) to see what each page contains:
     - extract_pdf_page(page=2)   → see what is on page 2
     - extract_pdf_page(page=3)   → see what is on page 3
     - extract_pdf_page(page=5)   → see what is on page 5
     The tool returns the output filename and image dimensions. Use these full-page PNGs directly.

  Step 0b — Crop ONLY if a page has a lot of white space above/below the figure AND the figure
     is clearly in a small region. The extraction description tells you the image size in pixels;
     if the figure is at the top of the page, crop with y0=0.0, y1=0.3 (for example).
     If unsure, use the FULL page — `\\includegraphics[width=0.9\\textwidth]` will scale it fine.
     DO NOT invent crop values like y0=0.3 — that is just a guess and will often be wrong.

  b) Reference extracted figures in LaTeX frames using \\figcap{} (NOT \\caption*):
     % Add to preamble: \\newcommand{\\figcap}[1]{\\par\\vspace{2pt}{\\scriptsize\\textit{#1}}}

     \\begin{frame}{Method Overview}
       \\centering
       \\includegraphics[width=0.9\\textwidth]{figures/page_3.png}
       \\figcap{Figure 2: Architecture from paper.}
     \\end{frame}

     % Side-by-side figure + text:
     \\begin{frame}{Results}
       \\begin{columns}[T]
         \\begin{column}{0.55\\textwidth}
           \\includegraphics[width=\\textwidth]{figures/page_5.png}
         \\end{column}
         \\begin{column}{0.40\\textwidth}
           \\begin{itemize}
             \\item<1-> Key finding 1
             \\item<2-> Key finding 2
           \\end{itemize}
         \\end{column}
       \\end{columns}
     \\end{frame}

  c) Parse ALL numeric results — every table in the paper must appear as a LaTeX tabular in slides
  d) If institution has a logo: download_file to fetch it, then \\includegraphics[height=0.5cm]{logo.png}

"""

        return f"""Build a {demo_form} ({form_hint}) for: "{paper.title}"

Contribution: {analysis.contribution}
Demo type: {demo_type}
Core concepts: {analysis.interaction_pattern}
{figures_note}
PRIORITY ORDER:
1. Extract key figures and result tables from the paper PDF using extract_pdf_page (MANDATORY FIRST STEP)
2. Parse ALL numeric results (every table, every benchmark) from the paper text
3. Plan all slide content before writing any code
4. Build slides with REAL content:
   - Figures: embed extracted PDF PNGs with slide.shapes.add_picture('figures/page_N.png', ...)
   - Tables: use slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table — NEVER use rectangles for tables
   - Charts: use slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, ...) (or LINE_MARKERS for trends)
     and explicitly set chart title, legend font color, axis tick/title colors, and data labels.
5. Run build.py with execute_python — must produce presentation.pptx without errors

CRITICAL RULES:
- Every result table in the paper MUST appear as an add_table() styled table in the slides
- Every key figure/diagram MUST be embedded (extracted from PDF with extract_pdf_page)
- No placeholder content — every slide uses real paper data
- NEVER simulate tables with filled rectangles — always use add_table()
"""

    def get_polish_prompt(self, paper, analysis, demo_form, demo_type, generated_files):
        spec = FORM_SPECS.get(demo_form, {})
        main_file = spec.get("main_file", "demo.html")
        quality_bar = spec.get("quality_bar", "publication quality")

        if demo_form == "slides":
            return f"""QUALITY REVIEW for python-pptx presentation — generated: {', '.join(generated_files[:12])}

Step 1 — Read build.py and audit critical requirements:
  • Does every slide call set_dark_bg()?
  • Does every content slide have add_header_bar() with the accent color?
  • Are ALL numeric result tables implemented with slide.shapes.add_table(...).table?
    → If any table was built using filled rectangles/textboxes, REWRITE it using add_table()
  • Is there at least one add_chart(XL_CHART_TYPE.*) with real paper comparison numbers?
  • Chart readability on dark background:
    → chart title, legend text, axis titles, and tick labels use light font color
    → plot/chart area fill matches dark theme; major gridlines are subtle but visible
    → show data labels on key charts (`plot.has_data_labels = True`)
  • Are there >=12 slides covering the full story?
  • Are key figures embedded using slide.shapes.add_picture('figures/page_N.png', ...)?
    → If figures/ directory exists, make sure they are used in the slides

Step 2 — Content completeness:
  • Title slide: paper title, all authors, venue, year?
  • Is there a results slide with real benchmark numbers?
  • Are speaker notes (slide.notes_slide.notes_text_frame.text = '...') on every slide?
  • Does build.py end with prs.save('presentation.pptx')?

Step 3 — Execute and verify:
  • Run execute_python to run build.py — must produce presentation.pptx without errors
  • Fix any python-pptx API errors found

Target quality: {quality_bar}"""

        elif demo_form == "latex":
            return f"""QUALITY REVIEW for Theoretical Explainer (LaTeX/Beamer) — generated: {', '.join(generated_files[:12])}

Step 1 — Read presentation.tex and audit:
  • Is documentclass[aspectratio=169]{{beamer}} correct?
  • Does every frame with bullets use \\item<1-> incremental reveals?
  • Are all equations in correct LaTeX math environments?
  • Is the TikZ architecture diagram present?

Step 2 — Content completeness:
  • Are there >=12 frames covering the paper?
  • Is the results table using booktabs (\\toprule/\\midrule/\\bottomrule)?
  • Are all actual numeric results from the paper in the table?
  • Do build.sh and README.md exist with compile instructions?

Step 3 — Syntax check:
  • execute_python to count unbalanced braces in the .tex file
  • Fix any LaTeX syntax errors found

Target quality: {quality_bar}"""

        elif demo_form == "presentation":
            return f"""QUALITY REVIEW for reveal.js presentation — generated: {', '.join(generated_files[:12])}

Step 1 — Read demo.html and count:
  • Is reveal.js loaded from unpkg.com/reveal.js@5.2.1? If not, update the CDN URL.
  • Are there >=14 <section> slides? If fewer, add missing slides (ablation, comparison, demo).
  • Does every list slide use class="fragment"? Find any <li> WITHOUT fragment and add it.
  • Is there an <aside class="notes">...</aside> block inside EVERY content <section>?
    → If ANY slide is missing speaker notes, add them now. Notes should be 2-3 sentences.
  • Are there >=3 original inline SVG diagrams? If fewer, draw the missing ones.

Step 2 — Content completeness:
  • Title slide: paper title, authors, venue, year?
  • Results slide: real benchmark numbers from the paper in a styled <table>?
  • Conclusion slide: 3-5 takeaways matching the paper's actual contributions?

Step 3 — Math & diagrams:
  • Do all KaTeX expressions use $$ (block) or $ (inline)? Fix any \\(...\\) or \\[...\\] notation.
  • Are SVG diagrams sized correctly (width/height attributes set)?

Target quality: {quality_bar}"""

        else:
            return f"""QUALITY REVIEW for Theoretical Explainer — generated: {', '.join(generated_files[:12])}

Step 1 — Read {main_file}:
  • Is every math expression using correct syntax?
  • Are there at least 3 original diagrams?
  • Do reveals/animations show concepts progressively?

Step 2 — Content:
  • Is there a section for each: Motivation, Key Insight, Main Result, Example?
  • Are intuitions written in plain English (no jargon without definition)?
  • Is there a concrete numerical example with actual numbers?

Step 3 — Visual quality:
  • Dark theme applied correctly?
  • Typography consistent: Inter font, correct weights?
  • All hover/transition states implemented?

Fix anything that falls short. Target quality: {quality_bar}"""
